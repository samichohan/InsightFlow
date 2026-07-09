"""Agent 7 — Professional Report Generator (PDF + DOCX + PPTX)."""

import os
import uuid
from datetime import datetime

from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image,
    Table, TableStyle, PageBreak
)

from app.agents.base import BaseAgent, PipelineContext
from app.core.config import settings


# ─────────────────────────────────────────────────────────────────────────────
# PDF Report
# ─────────────────────────────────────────────────────────────────────────────
def generate_pdf(dataset_name: str, quality: dict, eda: dict,
                 insights: list, recommendations: list,
                 chart_paths: dict, forecast: dict = None) -> str:

    filename = f"report_{uuid.uuid4().hex[:8]}.pdf"
    path = os.path.join(settings.REPORT_DIR, filename)
    doc = SimpleDocTemplate(path, pagesize=A4,
                             topMargin=0.7*inch, bottomMargin=0.7*inch)
    styles = getSampleStyleSheet()
    h1 = styles["Heading1"]
    h2 = styles["Heading2"]
    body = styles["BodyText"]
    title_style = ParagraphStyle("Title", parent=styles["Title"], fontSize=22)

    elements = []

    # Cover
    elements.append(Paragraph("AI Data Analyst Pro", title_style))
    elements.append(Paragraph(f"Report: {dataset_name}", body))
    elements.append(Paragraph(f"Generated: {datetime.now().strftime('%d %B %Y, %I:%M %p')}", body))
    elements.append(Spacer(1, 20))

    # Executive Summary
    elements.append(Paragraph("1. Executive Summary", h1))
    elements.append(Paragraph(
        f"This report analyzes the dataset '{dataset_name}' containing "
        f"{quality.get('total_rows', 'N/A')} rows and {quality.get('total_columns', 'N/A')} columns. "
        f"Overall Data Quality Score: {quality.get('data_quality_score', 'N/A')}/100.",
        body
    ))
    elements.append(Spacer(1, 12))

    # Data Quality
    elements.append(Paragraph("2. Data Quality Report", h1))
    q_data = [["Metric", "Value"]] + [
        [k, str(v)] for k, v in [
            ("Total Rows", quality.get("total_rows")),
            ("Total Columns", quality.get("total_columns")),
            ("Duplicate Rows", quality.get("duplicate_rows")),
            ("Total Missing Cells", quality.get("total_missing_cells")),
            ("Data Quality Score", f"{quality.get('data_quality_score')}/100"),
        ]
    ]
    t = Table(q_data, colWidths=[3*inch, 3*inch])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0,0),(-1,0), colors.HexColor("#1e3a5f")),
        ("TEXTCOLOR", (0,0),(-1,0), colors.white),
        ("GRID", (0,0),(-1,-1), 0.5, colors.grey),
        ("FONTSIZE", (0,0),(-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1),(-1,-1), [colors.white, colors.HexColor("#f1f5f9")]),
    ]))
    elements.append(t)
    elements.append(Spacer(1, 12))

    for s in quality.get("suggestions", []):
        elements.append(Paragraph(f"• {s}", body))
    elements.append(Spacer(1, 12))

    # Charts
    elements.append(Paragraph("3. Visualizations", h1))
    for name, cpath in chart_paths.items():
        if cpath and os.path.exists(cpath):
            elements.append(Paragraph(name.replace("_"," ").title(), h2))
            elements.append(Image(cpath, width=5.5*inch, height=3.5*inch))
            elements.append(Spacer(1, 10))
    elements.append(PageBreak())

    # Insights
    elements.append(Paragraph("4. Business Insights", h1))
    for insight in (insights or ["No insights generated."]):
        elements.append(Paragraph(f"• {insight}", body))
    elements.append(Spacer(1, 12))

    # Forecast
    if forecast and "error" not in forecast:
        elements.append(Paragraph("5. Forecast Report", h1))
        elements.append(Paragraph(
            f"Model: {forecast.get('model_type')} | R²: {forecast.get('r2_score')} | "
            f"MAE: {forecast.get('mae')} | Accuracy: {forecast.get('accuracy_percent')}%", body
        ))
        f_data = [["Period", "Forecasted Value"]] + list(zip(
            forecast.get("forecasted_periods", []),
            forecast.get("forecasted_values", [])
        ))
        ft = Table(f_data, colWidths=[3*inch, 3*inch])
        ft.setStyle(TableStyle([
            ("BACKGROUND", (0,0),(-1,0), colors.HexColor("#16a34a")),
            ("TEXTCOLOR", (0,0),(-1,0), colors.white),
            ("GRID", (0,0),(-1,-1), 0.5, colors.grey),
            ("FONTSIZE", (0,0),(-1,-1), 9),
        ]))
        elements.append(ft)
        elements.append(Spacer(1, 12))

    # Recommendations
    elements.append(Paragraph("6. Business Recommendations", h1))
    for rec in (recommendations or ["No recommendations generated."]):
        elements.append(Paragraph(f"• {rec}", body))

    doc.build(elements)
    return path


# ─────────────────────────────────────────────────────────────────────────────
# DOCX Report
# ─────────────────────────────────────────────────────────────────────────────
def generate_docx(dataset_name: str, quality: dict, eda: dict,
                  insights: list, recommendations: list,
                  chart_paths: dict, forecast: dict = None) -> str:

    doc = Document()
    filename = f"report_{uuid.uuid4().hex[:8]}.docx"
    path = os.path.join(settings.REPORT_DIR, filename)

    # Title
    t = doc.add_heading("AI Data Analyst Pro — Report", 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Dataset: {dataset_name}").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Generated: {datetime.now().strftime('%d %B %Y')}").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()

    # Executive Summary
    doc.add_heading("1. Executive Summary", 1)
    doc.add_paragraph(
        f"This report analyzes '{dataset_name}' with {quality.get('total_rows')} rows "
        f"and {quality.get('total_columns')} columns. "
        f"Data Quality Score: {quality.get('data_quality_score')}/100."
    )

    # Data Quality
    doc.add_heading("2. Data Quality", 1)
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Light Grid Accent 1"
    tbl.rows[0].cells[0].text = "Metric"
    tbl.rows[0].cells[1].text = "Value"
    for metric, value in [
        ("Total Rows", quality.get("total_rows")),
        ("Columns", quality.get("total_columns")),
        ("Duplicates", quality.get("duplicate_rows")),
        ("Missing Cells", quality.get("total_missing_cells")),
        ("Quality Score", f"{quality.get('data_quality_score')}/100"),
    ]:
        r = tbl.add_row().cells
        r[0].text, r[1].text = str(metric), str(value)

    doc.add_heading("Suggestions", 2)
    for s in quality.get("suggestions", []):
        doc.add_paragraph(s, style="List Bullet")

    # Charts
    doc.add_heading("3. Visualizations", 1)
    for name, cpath in chart_paths.items():
        if cpath and os.path.exists(cpath):
            doc.add_heading(name.replace("_", " ").title(), 2)
            doc.add_picture(cpath, width=Inches(5.5))

    # Insights
    doc.add_heading("4. Business Insights", 1)
    for insight in (insights or []):
        doc.add_paragraph(insight, style="List Bullet")

    # Forecast
    if forecast and "error" not in forecast:
        doc.add_heading("5. Forecast Report", 1)
        doc.add_paragraph(
            f"Model: {forecast.get('model_type')} | R²: {forecast.get('r2_score')} | "
            f"MAE: {forecast.get('mae')} | Accuracy: {forecast.get('accuracy_percent')}%"
        )

    # Recommendations
    doc.add_heading("6. Recommendations", 1)
    for rec in (recommendations or []):
        doc.add_paragraph(rec, style="List Bullet")

    doc.save(path)
    return path


# ─────────────────────────────────────────────────────────────────────────────
# PPTX Report
# ─────────────────────────────────────────────────────────────────────────────
def generate_pptx(dataset_name: str, insights: list, recommendations: list) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    filename = f"report_{uuid.uuid4().hex[:8]}.pptx"
    path = os.path.join(settings.REPORT_DIR, filename)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Data Analyst Pro"
    slide.placeholders[1].text = f"Dataset: {dataset_name}\n{datetime.now().strftime('%d %B %Y')}"

    # Insights slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Insights"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, insight in enumerate((insights or [])[:6]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {insight[:150]}"
        p.font.size = Pt(14)

    # Recommendations slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommendations"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, rec in enumerate((recommendations or [])[:6]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {rec[:150]}"
        p.font.size = Pt(14)

    prs.save(path)
    return path


class ReportAgent(BaseAgent):
    name = "report"

    def run(self, context: PipelineContext) -> dict:
        dataset_name = context.get("dataset_name", "dataset")
        quality = context.get("data_cleaning") or {}
        eda = context.get("eda") or {}
        insights = context.get("insights") or ["Insights not generated."]
        recommendations = context.get("recommendations") or ["Recommendations not generated."]
        charts = context.get("static_charts") or {}
        forecast = context.get("forecasting")
        fmt = context.get("report_format", "pdf")

        if fmt == "docx":
            path = generate_docx(dataset_name, quality, eda, insights, recommendations, charts, forecast)
        elif fmt == "pptx":
            path = generate_pptx(dataset_name, insights, recommendations)
        else:
            path = generate_pdf(dataset_name, quality, eda, insights, recommendations, charts, forecast)

        return {"report_path": path, "format": fmt}
